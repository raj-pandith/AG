"""
My Doctor — Premium PPT Generator
Generates a modern, dark-themed PowerPoint presentation from the presentation content.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ===== COLOR PALETTE =====
COLORS = {
    'bg_dark': RGBColor(15, 12, 41),       # #0f0c29
    'bg_mid': RGBColor(26, 26, 46),         # #1a1a2e
    'accent_purple': RGBColor(129, 140, 248),  # #818cf8
    'accent_violet': RGBColor(139, 92, 246),   # #8b5cf6
    'accent_pink': RGBColor(192, 132, 252),    # #c084fc
    'accent_blue': RGBColor(59, 130, 246),     # #3b82f6
    'accent_green': RGBColor(34, 197, 94),     # #22c55e
    'accent_red': RGBColor(239, 68, 68),       # #ef4444
    'accent_orange': RGBColor(249, 115, 22),   # #f97316
    'accent_teal': RGBColor(20, 184, 166),     # #14b8a6
    'white': RGBColor(255, 255, 255),
    'white_80': RGBColor(255, 255, 255),
    'white_60': RGBColor(200, 200, 220),
    'white_40': RGBColor(160, 160, 190),
    'card_bg': RGBColor(30, 30, 55),
    'card_border': RGBColor(60, 60, 100),
}

# Slide dimensions (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ===== HELPER FUNCTIONS =====

def set_slide_bg(slide, color: RGBColor):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_gradient_bg(slide, color1: RGBColor, color2: RGBColor, color3: RGBColor = None):
    """Add a gradient background to a slide using XML."""
    bg = slide.background
    fill = bg.fill
    # Use gradient fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2
    if color3:
        # Add a third stop
        gs = fill.gradient_stops._add_stop()
        # We'll use a simpler approach - just two colors for now


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None, radius=None):
    """Add a rectangle shape to a slide."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()

    if radius:
        # Set rounded corners via XML
        xml = shape._element
        spPr = xml.find(qn('p:spPr'))
        if spPr is None:
            spPr = etree.SubElement(xml, qn('p:spPr'))
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', f'val {int(radius * 10000)}')

    return shape


def add_text_box(slide, left, top, width, height, text, font_size, bold=False,
                 color=None, align=PP_ALIGN.LEFT, italic=False, line_spacing=None):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing

    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color

    return txBox


def add_paragraph(tf, text, font_size, bold=False, color=None,
                  align=PP_ALIGN.LEFT, italic=False, space_before=None):
    """Add a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = space_before

    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return p


def add_gradient_rect(slide, left, top, width, height, color1, color2, radius=None):
    """Add a rectangle with gradient fill."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.gradient()
    shape.fill.gradient_stops[0].color.rgb = color1
    shape.fill.gradient_stops[1].color.rgb = color2
    shape.line.fill.background()
    if radius:
        xml = shape._element
        spPr = xml.find(qn('p:spPr'))
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', f'val {int(radius * 10000)}')
    return shape


def add_card(slide, left, top, width, height, icon, title, desc, icon_color=None, title_color=None):
    """Add a premium card with icon, title, and description."""
    # Card background
    card = add_rect(slide, left, top, width, height,
                    fill_color=COLORS['card_bg'],
                    line_color=COLORS['card_border'],
                    line_width=Pt(0.5),
                    radius=20000)

    # Top accent line
    accent = add_rect(slide, left + Inches(0.08), top + Inches(0.05),
                      width - Inches(0.16), Inches(0.04),
                      fill_color=COLORS['accent_purple'])

    # Icon
    icon_box = add_text_box(slide, left + Inches(0.15), top + Inches(0.14),
                            Inches(0.6), Inches(0.5), icon, Pt(22))

    # Title
    title_box = add_text_box(slide, left + Inches(0.15), top + Inches(0.42),
                             width - Inches(0.3), Inches(0.5),
                             title, Pt(16), bold=True, color=COLORS['accent_purple'])

    # Description
    desc_box = add_text_box(slide, left + Inches(0.15), top + Inches(0.65),
                            width - Inches(0.3), height - Inches(0.85),
                            desc, Pt(12), color=COLORS['white_60'])


def add_numbered_item(slide, left, top, width, num, text, num_bg, num_text_color=COLORS['white'], text_color=None):
    """Add a numbered list item."""
    item_h = Inches(0.65)
    # Background
    bg = add_rect(slide, left, top, width, item_h,
                  fill_color=RGBColor(25, 25, 50),
                  line_color=COLORS['card_border'],
                  line_width=Pt(0.5),
                  radius=15000)

    # Number circle
    circle_size = Inches(0.38)
    circle = add_rect(slide, left + Inches(0.12), top + Inches(0.13),
                      circle_size, circle_size,
                      fill_color=num_bg,
                      radius=50000)
    num_box = add_text_box(slide, left + Inches(0.12), top + Inches(0.13),
                           circle_size, circle_size,
                           str(num), Pt(16), bold=True,
                           color=num_text_color, align=PP_ALIGN.CENTER)

    # Text
    text_box = add_text_box(slide, left + Inches(0.62), top + Inches(0.13),
                            width - Inches(0.75), item_h - Inches(0.1),
                            text, Pt(13), color=text_color or COLORS['white_80'])


def add_section_tag(slide, text, left, top):
    """Add a small section tag/label."""
    tag_w = Inches(2.2)
    tag_h = Inches(0.32)
    tag = add_rect(slide, left, top, tag_w, tag_h,
                   fill_color=RGBColor(35, 35, 65),
                   line_color=COLORS['accent_purple'],
                   line_width=Pt(0.5),
                   radius=50000)
    tag_box = add_text_box(slide, left, top, tag_w, tag_h,
                           text, Pt(10), bold=True,
                           color=COLORS['accent_purple'],
                           align=PP_ALIGN.CENTER)


def add_dot_indicator(slide, left, top):
    """Add a small pulsing dot indicator."""
    dot = slide.shapes.add_shape(9, left, top, Inches(0.12), Inches(0.12))  # Oval
    dot.fill.solid()
    dot.fill.fore_color.rgb = COLORS['accent_purple']
    dot.line.fill.background()
    return dot


# ===== SLIDE BUILDERS =====

def slide_hero(prs):
    """Slide 1: Hero - My Doctor"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, COLORS['bg_dark'])

    # Decorative circles
    circ1 = slide.shapes.add_shape(9, Inches(9.5), Inches(-1.5), Inches(4), Inches(4))
    circ1.fill.solid()
    circ1.fill.fore_color.rgb = RGBColor(99, 102, 241)
    circ1.fill.fore_color.brightness = 0.6
    circ1.line.fill.background()

    circ2 = slide.shapes.add_shape(9, Inches(-1.5), Inches(5), Inches(3), Inches(3))
    circ2.fill.solid()
    circ2.fill.fore_color.rgb = RGBColor(139, 92, 246)
    circ2.fill.fore_color.brightness = 0.6
    circ2.line.fill.background()

    # Main Icon
    add_text_box(slide, Inches(5.5), Inches(1.5), Inches(2.5), Inches(1.2),
                 "🩺", Pt(80), align=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10.3), Inches(1.3),
                 "My Doctor", Pt(72), bold=True, color=COLORS['white'],
                 align=PP_ALIGN.CENTER)

    # Underline accent
    line = add_rect(slide, Inches(4.5), Inches(4.2), Inches(4.3), Inches(0.06),
                    fill_color=COLORS['accent_purple'])

    # Subtitle
    add_text_box(slide, Inches(1.5), Inches(4.4), Inches(10.3), Inches(0.8),
                 "Your Digital Healthcare Companion", Pt(24), color=COLORS['white_60'],
                 align=PP_ALIGN.CENTER)

    # Badge
    badge = add_rect(slide, Inches(4.8), Inches(5.4), Inches(3.7), Inches(0.45),
                     fill_color=RGBColor(35, 35, 70),
                     line_color=RGBColor(80, 80, 140),
                     line_width=Pt(0.5),
                     radius=50000)
    add_text_box(slide, Inches(4.8), Inches(5.4), Inches(3.7), Inches(0.45),
                 "●  Connecting Patients & Doctors Seamlessly", Pt(13),
                 color=COLORS['accent_purple'], align=PP_ALIGN.CENTER)


def slide_vision(prs):
    """Slide 2: Vision"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(10, 10, 35))

    # Decorative
    circ = slide.shapes.add_shape(9, Inches(10.5), Inches(4.5), Inches(3.5), Inches(3.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = RGBColor(99, 102, 241)
    circ.fill.fore_color.brightness = 0.5
    circ.line.fill.background()

    # Tag
    add_section_tag(slide, "OUR VISION", Inches(5.5), Inches(1.3))

    # Icon
    add_text_box(slide, Inches(6), Inches(1.75), Inches(1.3), Inches(1),
                 "🌟", Pt(50), align=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, Inches(1.5), Inches(2.9), Inches(10.3), Inches(0.8),
                 "Our Vision", Pt(44), bold=True, color=COLORS['white'],
                 align=PP_ALIGN.CENTER)

    # Vision text in a card
    text_card = add_rect(slide, Inches(1.5), Inches(4.0), Inches(10.3), Inches(2.5),
                         fill_color=RGBColor(25, 25, 55),
                         line_color=RGBColor(70, 70, 120),
                         line_width=Pt(0.5),
                         radius=20000)

    vision_text = ('"To become the digital healthcare companion that keeps patients '
                   'connected with their trusted doctors while creating a unified '
                   'medical record that improves healthcare for everyone."')

    tf = slide.shapes.add_textbox(Inches(2), Inches(4.3), Inches(9.3), Inches(2.0)).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    # "digital healthcare companion" in purple
    run = p.add_run()
    run.text = "To become the "
    run.font.size = Pt(22)
    run.font.color.rgb = COLORS['white_80']

    run2 = p.add_run()
    run2.text = "digital healthcare companion"
    run2.font.size = Pt(22)
    run2.font.bold = True
    run2.font.color.rgb = COLORS['accent_purple']

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)
    run3 = p2.add_run()
    run3.text = "that keeps patients connected with their trusted doctors while "
    run3.font.size = Pt(22)
    run3.font.color.rgb = COLORS['white_80']

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(4)
    run4 = p3.add_run()
    run4.text = "creating a unified medical record that "
    run4.font.size = Pt(22)
    run4.font.color.rgb = COLORS['white_80']

    run5 = p3.add_run()
    run5.text = "improves healthcare for everyone"
    run5.font.size = Pt(22)
    run5.font.bold = True
    run5.font.color.rgb = COLORS['accent_purple']

    run6 = p3.add_run()
    run6.text = '."'
    run6.font.size = Pt(22)
    run6.font.color.rgb = COLORS['white_80']


def slide_problems_overview(prs):
    """Slide 3: Problems Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(12, 15, 40))

    add_section_tag(slide, "THE CHALLENGE", Inches(5.5), Inches(0.9))
    add_text_box(slide, Inches(1.5), Inches(1.4), Inches(10.3), Inches(0.8),
                 "Healthcare Has Real Problems", Pt(40), bold=True,
                 color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(2.3), Inches(10.3), Inches(0.6),
                 ("While technology has advanced, the fundamental connection between "
                  "patients and doctors is broken."),
                 Pt(16), color=COLORS['white_60'], align=PP_ALIGN.CENTER)

    # Doctor side
    doc_card = add_rect(slide, Inches(1.8), Inches(3.2), Inches(4.2), Inches(3.2),
                        fill_color=RGBColor(20, 20, 50),
                        line_color=COLORS['accent_violet'],
                        line_width=Pt(1),
                        radius=20000)
    add_text_box(slide, Inches(2.0), Inches(3.4), Inches(4), Inches(0.7),
                 "👨‍⚕️", Pt(36), align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.0), Inches(4.1), Inches(4), Inches(0.5),
                 "Doctors", Pt(20), bold=True, color=COLORS['accent_purple'],
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.0), Inches(4.6), Inches(4), Inches(1.2),
                 ("Lose patients when away\n"
                  "Lack patient context\n"
                  "No continuity of care"),
                 Pt(13), color=COLORS['white_60'], align=PP_ALIGN.CENTER)

    # Arrow
    add_text_box(slide, Inches(6.1), Inches(4.0), Inches(1.1), Inches(0.8),
                 "⟶", Pt(40), color=COLORS['white_40'], align=PP_ALIGN.CENTER)

    # Patient side
    pat_card = add_rect(slide, Inches(7.3), Inches(3.2), Inches(4.2), Inches(3.2),
                        fill_color=RGBColor(20, 20, 50),
                        line_color=COLORS['accent_red'],
                        line_width=Pt(1),
                        radius=20000)
    add_text_box(slide, Inches(7.5), Inches(3.4), Inches(4), Inches(0.7),
                 "👤", Pt(36), align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.5), Inches(4.1), Inches(4), Inches(0.5),
                 "Patients", Pt(20), bold=True, color=COLORS['accent_red'],
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.5), Inches(4.6), Inches(4), Inches(1.2),
                 ("Wait long hours\n"
                  "Carry paper files\n"
                  "No central records"),
                 Pt(13), color=COLORS['white_60'], align=PP_ALIGN.CENTER)


def slide_doctor_problems(prs):
    """Slide 4: Doctor Problems"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(16, 16, 38))

    add_section_tag(slide, "👨‍⚕️ DOCTOR'S SIDE", Inches(5.0), Inches(0.7))
    add_text_box(slide, Inches(1.5), Inches(1.2), Inches(10.3), Inches(0.8),
                 "Problems Doctors Face Every Day", Pt(36), bold=True,
                 color=COLORS['white'], align=PP_ALIGN.CENTER)

    problems = [
        ("Patient issues are closely related to doctors, but not all problems are communicated well.",
         COLORS['accent_red']),
        ("Doctors cannot lose patients even when they move out of town.",
         COLORS['accent_orange']),
        ("In the doctor's absence, AI can check patient history and suggest options — same medicine, referral, or availability info.",
         COLORS['accent_blue']),
        ("During vacation, if a doctor gets free time, they can attend patients online.",
         COLORS['accent_green']),
    ]

    for i, (text, bg_color) in enumerate(problems):
        top = Inches(2.2) + i * Inches(1.1)
        add_numbered_item(slide, Inches(1.5), top, Inches(10.3), i + 1, text, bg_color)


def slide_patient_problems(prs):
    """Slide 5: Patient's Problems"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(12, 14, 32))

    add_section_tag(slide, "👤 PATIENT'S SIDE", Inches(5.2), Inches(0.7))
    add_text_box(slide, Inches(1.5), Inches(1.2), Inches(10.3), Inches(0.8),
                 "Problems Patients Face Every Day", Pt(36), bold=True,
                 color=COLORS['white'], align=PP_ALIGN.CENTER)

    problems = [
        ("Online consultations with unknown doctors fail to understand patient issues properly.",
         COLORS['accent_red']),
        ("Patients have to wait in hospitals for long hours to get treated.",
         COLORS['accent_orange']),
        ("Token systems exist but lack live updates on the current status.",
         COLORS['accent_blue']),
        ("Sometimes patients can't go to the hospital and no family member is available to help.",
         COLORS['accent_green']),
        ("Patients must carry hard copies of medical reports and X-rays everywhere.",
         COLORS['accent_violet']),
        ("There is no centralized patient medical history system.",
         COLORS['accent_teal']),
    ]

    for i, (text, bg_color) in enumerate(problems):
        top = Inches(2.2) + i * Inches(0.82)
        add_numbered_item(slide, Inches(1.5), top, Inches(10.3), i + 1, text, bg_color)


def slide_patient_solutions_list(prs):
    """Slide 6: Patient Solutions (List view)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(14, 20, 45))

    add_section_tag(slide, "✅ PATIENT SOLUTIONS", Inches(5.0), Inches(0.7))
    add_text_box(slide, Inches(1.5), Inches(1.2), Inches(10.3), Inches(0.8),
                 "Our Smart Solutions", Pt(36), bold=True,
                 color=COLORS['white'], align=PP_ALIGN.CENTER)

    solutions = [
        ("Connect with a known doctor — they understand your history. "
         "A secure platform replaces the need to share personal numbers.",
         COLORS['accent_purple']),
        ("Live hospital updates — real-time doctor availability with "
         "GPS-based smart alerts to reach the hospital just in time.",
         COLORS['accent_blue']),
        ("Family member support — add family members who can set up calls "
         "and explain the situation on your behalf.",
         COLORS['accent_green']),
        ("Centralized cloud records — all medical history securely stored. "
         "Full summary sent to doctor on appointment. No more carrying files!",
         COLORS['accent_teal']),
    ]

    for i, (text, bg_color) in enumerate(solutions):
        top = Inches(2.2) + i * Inches(1.15)
        # Green checkmark circle
        check = add_rect(slide, Inches(1.5), top + Inches(0.1), Inches(0.5), Inches(0.5),
                         fill_color=bg_color, radius=50000)
        add_text_box(slide, Inches(1.5), top + Inches(0.1), Inches(0.5), Inches(0.5),
                     "✓", Pt(18), bold=True, color=COLORS['white'],
                     align=PP_ALIGN.CENTER)

        # Text bg
        bg = add_rect(slide, Inches(2.15), top, Inches(9.65), Inches(0.85),
                      fill_color=RGBColor(20, 25, 55),
                      line_color=RGBColor(50, 50, 90),
                      line_width=Pt(0.5),
                      radius=15000)
        add_text_box(slide, Inches(2.3), top + Inches(0.12), Inches(9.35), Inches(0.65),
                     text, Pt(13), color=COLORS['white_80'])


def slide_doctor_solutions_cards(prs):
    """Slide 7: Doctor Solutions - Cards"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(16, 16, 42))

    add_section_tag(slide, "🤖 AI-POWERED DOCTOR SUPPORT", Inches(4.8), Inches(0.6))
    add_text_box(slide, Inches(1.5), Inches(1.1), Inches(10.3), Inches(0.7),
                 "How We Help Doctors", Pt(36), bold=True,
                 color=COLORS['white'], align=PP_ALIGN.CENTER)

    cards = [
        ("🧠", "AI Patient Assistant",
         "When unavailable, AI reviews patient history and suggests the best next step — "
         "same medicine, referral, or appointment rescheduling."),
        ("🏖️", "Vacation Mode",
         "Doctors on vacation can still handle urgent online consultations in their free time. "
         "No patient is left unattended."),
        ("📋", "Smart Patient Handover",
         "Before any consultation, the system auto-sends a summarized patient history to "
         "the doctor — complete context in seconds."),
        ("🔒", "Secure Connection",
         "Doctors connect through a secure platform. No need to share personal phone numbers "
         "with every patient."),
    ]

    card_w = Inches(2.9)
    card_h = Inches(3.5)
    gap = Inches(0.22)
    total_w = 4 * card_w + 3 * gap
    start_x = (SLIDE_W - total_w) / 2

    for i, (icon, title, desc) in enumerate(cards):
        left = start_x + i * (card_w + gap)
        add_card(slide, left, Inches(2.0), card_w, card_h, icon, title, desc)


def slide_patient_solutions_cards(prs):
    """Slide 8: Patient Solutions - Cards"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(10, 14, 35))

    add_section_tag(slide, "📱 PATIENT EXPERIENCE REDEFINED", Inches(4.5), Inches(0.6))
    add_text_box(slide, Inches(1.5), Inches(1.1), Inches(10.3), Inches(0.7),
                 "A Better Way to Healthcare", Pt(36), bold=True,
                 color=COLORS['white'], align=PP_ALIGN.CENTER)

    cards = [
        ("👨‍⚕️", "Trusted Doctors",
         "Connect with doctors who already know your history. "
         "Better conversations, better care."),
        ("⏱️", "No More Waiting",
         "GPS-based smart alerts tell you exactly when to reach the hospital. "
         "Zero waiting time."),
        ("👨‍👩‍👧", "Family Support",
         "Add family members who can represent you during consultations. "
         "Great for elders and children."),
        ("☁️", "Cloud Records",
         "All reports, X-rays & history stored securely in the cloud. "
         "Access from anywhere, anytime."),
    ]

    card_w = Inches(2.9)
    card_h = Inches(3.5)
    gap = Inches(0.22)
    total_w = 4 * card_w + 3 * gap
    start_x = (SLIDE_W - total_w) / 2

    for i, (icon, title, desc) in enumerate(cards):
        left = start_x + i * (card_w + gap)
        add_card(slide, left, Inches(2.0), card_w, card_h, icon, title, desc)


def slide_how_it_works(prs):
    """Slide 9: How It Works - Roadmap"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(18, 18, 38))

    add_section_tag(slide, "HOW IT WORKS", Inches(5.7), Inches(0.8))
    add_text_box(slide, Inches(1.5), Inches(1.35), Inches(10.3), Inches(0.8),
                 "Seamless Experience, Every Step", Pt(38), bold=True,
                 color=COLORS['white'], align=PP_ALIGN.CENTER)

    steps = [
        ("1", "Book", "Book appointment with\nyour trusted doctor", COLORS['accent_purple']),
        ("→", "Get Summary", "Full medical history\nsent to doctor", COLORS['accent_violet']),
        ("3", "Consult", "Visit or consult online\nwith full context", COLORS['accent_pink']),
        ("4", "Track", "Records updated &\navailable anywhere", COLORS['accent_teal']),
    ]

    step_w = Inches(2.5)
    step_h = Inches(3.5)
    gap = Inches(0.4)
    total_w = 4 * step_w + 3 * gap
    start_x = (SLIDE_W - total_w) / 2

    for i, (num, label, desc, accent) in enumerate(steps):
        left = start_x + i * (step_w + gap)

        # Circle
        circle = slide.shapes.add_shape(9, left + Inches(0.85), Inches(2.2),
                                        Inches(0.8), Inches(0.8))
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent
        circle.line.fill.background()
        add_text_box(slide, left + Inches(0.85), Inches(2.2), Inches(0.8), Inches(0.8),
                     num, Pt(28), bold=True, color=COLORS['white'],
                     align=PP_ALIGN.CENTER)

        # Label
        add_text_box(slide, left, Inches(3.15), step_w, Inches(0.5),
                     label, Pt(18), bold=True, color=COLORS['white'],
                     align=PP_ALIGN.CENTER)

        # Description
        add_text_box(slide, left, Inches(3.6), step_w, Inches(1.0),
                     desc, Pt(12), color=COLORS['white_60'],
                     align=PP_ALIGN.CENTER)

        # Arrow (except last)
        if i < 3:
            arrow_x = left + step_w + Inches(0.02)
            add_text_box(slide, arrow_x, Inches(2.4), Inches(0.36), Inches(0.6),
                         "→", Pt(28), color=COLORS['white_40'], align=PP_ALIGN.CENTER)


def slide_common_ecosystem(prs):
    """Slide 10: Common / Ecosystem"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(8, 18, 25))

    add_section_tag(slide, "🌐 BEYOND HEALTHCARE", Inches(5.3), Inches(0.6))
    add_text_box(slide, Inches(1.5), Inches(1.15), Inches(10.3), Inches(0.7),
                 "Building a Complete Health Ecosystem", Pt(36), bold=True,
                 color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(1.9), Inches(10.3), Inches(0.5),
                 ("With centralized medical history, we unlock powerful features "
                  "that benefit everyone."),
                 Pt(14), color=COLORS['white_60'], align=PP_ALIGN.CENTER)

    features_row1 = [
        ("🛡️", "Better Insurance", "Accurate history = better insurance plans & faster claims."),
        ("⚡", "Fast Processing", "Digital records speed up approvals & reimbursements."),
        ("💊", "Generic Medicine", "Smart suggestions for affordable generic alternatives."),
    ]

    features_row2 = [
        ("⌚", "Wearables", "Sync with smart watches for real-time health monitoring."),
        ("👥", "Health Community", "Engage with a community for tips, support & motivation."),
        ("🚀", "And More...", "Endless possibilities with a unified health platform."),
    ]

    # Row 1
    card_w = Inches(3.8)
    card_h = Inches(2.0)
    gap_x = Inches(0.2)
    total_w1 = 3 * card_w + 2 * gap_x
    start_x = (SLIDE_W - total_w1) / 2

    for i, (icon, title, desc) in enumerate(features_row1):
        left = start_x + i * (card_w + gap_x)
        add_card(slide, left, Inches(2.6), card_w, card_h, icon, title, desc)

    # Row 2
    start_y = Inches(4.85)
    for i, (icon, title, desc) in enumerate(features_row2):
        left = start_x + i * (card_w + gap_x)
        add_card(slide, left, start_y, card_w, card_h, icon, title, desc)


def slide_impact(prs):
    """Slide 11: Impact Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(12, 14, 38))

    add_section_tag(slide, "IMPACT", Inches(5.8), Inches(0.7))
    add_text_box(slide, Inches(1.5), Inches(1.3), Inches(10.3), Inches(0.8),
                 "A Win-Win for Everyone", Pt(40), bold=True,
                 color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Three columns
    cols = [
        ("👨‍⚕️", "For Doctors", COLORS['accent_purple'],
         ["Never lose a patient",
          "Work smarter with AI",
          "Stay connected on vacation"]),
        ("👤", "For Patients", COLORS['accent_red'],
         ["No waiting in hospitals",
          "No paper files",
          "Doctors who know you"]),
        ("🏥", "For Healthcare", COLORS['accent_green'],
         ["Better insurance",
          "Faster processing",
          "Healthier communities"]),
    ]

    col_w = Inches(3.6)
    col_gap = Inches(0.4)
    total_w = 3 * col_w + 2 * col_gap
    start_x = (SLIDE_W - total_w) / 2

    for i, (icon, label, accent, bullets) in enumerate(cols):
        left = start_x + i * (col_w + col_gap)

        # Card
        card = add_rect(slide, left, Inches(2.3), col_w, Inches(4.0),
                        fill_color=RGBColor(22, 22, 50),
                        line_color=accent,
                        line_width=Pt(1),
                        radius=20000)

        # Icon
        add_text_box(slide, left, Inches(2.55), col_w, Inches(0.8),
                     icon, Pt(40), align=PP_ALIGN.CENTER)

        # Title
        add_text_box(slide, left, Inches(3.35), col_w, Inches(0.5),
                     label, Pt(18), bold=True, color=accent,
                     align=PP_ALIGN.CENTER)

        # Bullets
        for j, bullet in enumerate(bullets):
            add_text_box(slide, left + Inches(0.3), Inches(4.05) + j * Inches(0.45),
                         col_w - Inches(0.6), Inches(0.4),
                         f"• {bullet}", Pt(12), color=COLORS['white_60'])


def slide_roadmap(prs):
    """Slide 12: Future Roadmap"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(16, 14, 28))

    add_section_tag(slide, "ROADMAP", Inches(5.8), Inches(0.7))
    add_text_box(slide, Inches(1.5), Inches(1.3), Inches(10.3), Inches(0.8),
                 "What's Next for My Doctor", Pt(40), bold=True,
                 color=COLORS['white'], align=PP_ALIGN.CENTER)

    phases = [
        ("📱", "MVP Launch", "Core platform with appointment\nbooking & medical records",
         COLORS['accent_purple']),
        ("🧠", "AI Integration", "Smart AI assistant for\ndoctors & patients",
         COLORS['accent_violet']),
        ("⌚", "Wearables", "Real-time health monitoring\n& smart alerts",
         COLORS['accent_pink']),
        ("🌐", "Full Ecosystem", "Insurance, pharmacy &\ncommunity integration",
         COLORS['accent_teal']),
    ]

    phase_w = Inches(2.6)
    gap = Inches(0.3)
    total_w = 4 * phase_w + 3 * gap
    start_x = (SLIDE_W - total_w) / 2
    start_y = Inches(2.5)

    for i, (icon, label, desc, accent) in enumerate(phases):
        left = start_x + i * (phase_w + gap)

        # Circle
        circle = slide.shapes.add_shape(9, left + Inches(0.75), start_y,
                                        Inches(1.1), Inches(1.1))
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent
        circle.line.fill.background()
        add_text_box(slide, left + Inches(0.75), start_y, Inches(1.1), Inches(1.1),
                     icon, Pt(32), align=PP_ALIGN.CENTER)

        # Label
        add_text_box(slide, left, start_y + Inches(1.3), phase_w, Inches(0.5),
                     label, Pt(16), bold=True, color=COLORS['white'],
                     align=PP_ALIGN.CENTER)

        # Desc
        add_text_box(slide, left, start_y + Inches(1.75), phase_w, Inches(1.0),
                     desc, Pt(12), color=COLORS['white_60'],
                     align=PP_ALIGN.CENTER)

        # Arrow (except last)
        if i < 3:
            arrow_x = left + phase_w + Inches(0.02)
            add_text_box(slide, arrow_x, start_y + Inches(0.3), Inches(0.26), Inches(0.6),
                         "→", Pt(26), color=COLORS['white_40'], align=PP_ALIGN.CENTER)


def slide_closing(prs):
    """Slide 13: Closing"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS['bg_dark'])

    # Decorative
    circ1 = slide.shapes.add_shape(9, Inches(10), Inches(-2), Inches(4), Inches(4))
    circ1.fill.solid()
    circ1.fill.fore_color.rgb = RGBColor(99, 102, 241)
    circ1.fill.fore_color.brightness = 0.5
    circ1.line.fill.background()

    circ2 = slide.shapes.add_shape(9, Inches(-2), Inches(5), Inches(3.5), Inches(3.5))
    circ2.fill.solid()
    circ2.fill.fore_color.rgb = RGBColor(139, 92, 246)
    circ2.fill.fore_color.brightness = 0.5
    circ2.line.fill.background()

    # Icon
    add_text_box(slide, Inches(5.5), Inches(1.5), Inches(2.5), Inches(1.2),
                 "🩺", Pt(72), align=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10.3), Inches(1.0),
                 "My Doctor", Pt(60), bold=True, color=COLORS['white'],
                 align=PP_ALIGN.CENTER)

    # Accent line
    line = add_rect(slide, Inches(4.5), Inches(3.95), Inches(4.3), Inches(0.06),
                    fill_color=COLORS['accent_purple'])

    # Tagline
    tagline_tf = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.8)).text_frame
    tagline_tf.word_wrap = True
    p = tagline_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    r1 = p.add_run()
    r1.text = "Healthcare, reimagined — "
    r1.font.size = Pt(22)
    r1.font.color.rgb = COLORS['white_60']

    r2 = p.add_run()
    r2.text = "connected, intelligent, and personal."
    r2.font.size = Pt(22)
    r2.font.bold = True
    r2.font.color.rgb = COLORS['accent_purple']


# ===== MAIN =====

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Generating slides...")

    slides_builders = [
        slide_hero,
        slide_vision,
        slide_problems_overview,
        slide_doctor_problems,
        slide_patient_problems,
        slide_patient_solutions_list,
        slide_doctor_solutions_cards,
        slide_patient_solutions_cards,
        slide_how_it_works,
        slide_common_ecosystem,
        slide_impact,
        slide_roadmap,
        slide_closing,
    ]

    for builder in slides_builders:
        builder(prs)
        print(f"  -> {builder.__doc__.split(':')[0] if builder.__doc__ else builder.__name__}")

    output_path = "c:/My Web Sites/Jain2/my-doctor-presentation.pptx"
    prs.save(output_path)
    print(f"\nDone! Saved: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
